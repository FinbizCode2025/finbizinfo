from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def add_slide(title, content_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            if item.startswith("-"):
                 p.text = item[1:].strip()
                 p.level = 1

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "FinBiz Info"
    subtitle.text = "Financial Analysis & Intelligence Platform\nAutomated Extraction, Analysis, and AI Insights"

    # Slide 2: Key Features
    add_slide("Key Features", [
        "Financial Document Processing: High-performance PDF extraction.",
        "Automated Ratio Engine: 40+ Financial Ratios calculated instantly.",
        "AI Multi-Agent System: Orchestrated agents for validation & quality.",
        "ERP Integration: Direct connection to Tally ERP.",
        "Corporate Data: Integration with Probe42 API."
    ])

    # Slide 3: Financial Document Processing
    add_slide("Document Processing", [
        "PyMuPDF Extraction: High-speed text and layout parsing.",
        "OCR Support: Processing scanned documents via GLM/Ollama.",
        "Table Reconstruction: Vertical & Horizontal clustering heuristics.",
        "Intelligent Note Selection: Targeting key account notes for metrics."
    ])

    # Slide 4: Automated Financial Analysis
    add_slide("Financial Ratio Engine", [
        "Liquidity: Current, Quick, Cash Ratios.",
        "Profitability: Margins (Gross, Net, EBITDA), ROA, ROE.",
        "Solvency: Debt-to-Equity, Interest Coverage.",
        "Efficiency: Asset Turnover, Receivables Turnover, DSO.",
        "DuPont Analysis: 3-step breakdown of ROE."
    ])

    # Slide 5: AI & Multi-Agent Orchestration
    add_slide("Multi-Agent Intelligence", [
        "FinancialRatioAgent: Calculation & interpretation tasks.",
        "DataValidationAgent: Internal consistency (Assets = E + L).",
        "QualityCheckAgent: Reliability assessment.",
        "FactCheckingAgent: PDF source citation and verification.",
        "RAG: Natural language querying over documents using LangChain."
    ])

    # Slide 6: Technology Stack
    add_slide("Technology Stack", [
        "Backend: Python, Flask, LangChain, LangGraph.",
        "AI/LLM: Google Gemini, Kimi (Moonshot), Mistral.",
        "Data: FAISS (Vector DB), Pandas, NumPy, SQLAlchemy.",
        "Frontend: React, TypeScript, Vite, Tailwind CSS.",
        "External: Tavily (Search), Probe42, Tally ERP."
    ])

    # Slide 7: Project Structure
    add_slide("Project Structure", [
        "/be: Backend (app.py, agents.py, analyzer.py, models.py).",
        "/fe: Frontend (React components, Pages, Utils).",
        "/uploads: Session and PDF storage.",
        "finbiz.db: SQLite database."
    ])

    # Slide 8: Logic Flow
    add_slide("Core Logic Flow", [
        "1. Extraction: PDF parsing and OCR.",
        "2. Refinement: LLM-based structural mapping (Kimi/Gemini).",
        "3. Intelligence: Agent execution and ratio calculation.",
        "4. Interaction: RAG-based chat and report generation."
    ])

    prs.save('FinBiz_Info_Documentation.pptx')
    print("Presentation saved as FinBiz_Info_Documentation.pptx")

if __name__ == "__main__":
    create_presentation()
